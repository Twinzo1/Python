# -*- coding: utf-8 -*-
"""
@Time ： 2021/6/4 21:43
@Auth ： Twinzo1
@File ：bdwenku.py
@IDE ：PyCharm
@Motto：ABC(Always Be Coding)
@Version: V3.0
@Description: 百度文库爬取
依赖库: requests, python_docx, python_pptx, pyquery
自动识别下载文件格式
"""
import requests
import json
import re
import os
from pyquery import PyQuery as pq
from docx import Document
from docx.shared import Inches, Pt
from docx.enum.text import WD_LINE_SPACING
import pptx
from pptx.util import Inches

def str2dict(dict_str):
    """
    等号分号(= ;)型字符串转字典
    :return: 字典cookie
    """
    json_data = json.dumps(dict_str)
    json_str = json.loads(json_data)
    ret_dict = dict([item.split("=", 1) for item in json_str.split(";")])
    return ret_dict

def format_dict_str(dict_str):
    """
    用于格式化不符合格式的字典字符串，如{"key":"value "EEE""},将之格式化为{"key":"value \"EEE\""}
    :param dict_str:
    :return:
    """
    first_split = re.split(':', dict_str)
    colon_str = []
    for sp1 in first_split:
        second_split = re.split(",", sp1)
        comma_str = []
        for sp2 in second_split:
            if sp2.count('"') == 2 or sp2.count('"') == 0:
                comma_str.append(sp2)
                continue
            if sp2.count('"') - 2 != sp2.count(r'\"'):
                splitBack = sp2.replace(r'"', r'\"')
                if sp2.count('"') % 2 != 0:
                    # 替换最后一个\"
                    sp2 = splitBack[::-1].replace('"\\', '"', 1)[::-1]
                    comma_str.append(sp2)
                    continue
                splitBack1 = splitBack.replace(r'\"', '"', 1)
                # 替换最后一个\"
                sp2 = splitBack1[::-1].replace('"\\', '"', 1)[::-1]
            comma_str.append(sp2)
        sp1 = ",".join(comma_str)
        colon_str.append(sp1)
    return ":".join(colon_str)


class BDWenKu:
    def __init__(self, bd_url, bd_cookie=""):
        self.headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) "
                          "Chrome/78.0.3904.116 Safari/537.36 "
        }
        self.cookie = str2dict(bd_cookie) if bd_cookie != "" else ""
        self.url = bd_url
        self.pq_html = self.get_html()
        self.title = self.get_title()
        self.homepage_script_dict = self.get_homepage_content_dict()
        self.file_type = self.get_file_type()

    def get_html(self):
        """
        获取网页html
        :return:
        """
        headers = {"user-agent": "LogStatistic"}
        res = self.with_cookie(self.url, headers)
        doc = pq(res.text)
        return doc

    def get_title(self):
        title = self.pq_html('title')
        return title.text().replace(" - 百度文库", "")

    def with_cookie(self, url, headers=""):
        headers = self.headers if headers == "" else headers
        if not self.cookie or len(str(self.cookie).strip()) == 0:
            res = requests.get(url, headers=headers)
        else:
            res = requests.get(url, cookies=self.cookie, headers=headers)
        res.encoding = res.apparent_encoding
        return res

    def get_homepage_content_dict(self):
        """
        获取主页的数据
        :return:返回字典类型
        """
        content = self.pq_html('script')
        for url_info in content.items():
            if "var pageData" in url_info.text():
                url_list = re.findall(r"=(.+);.?window.pageData", url_info.text())
                content_have_url_dict = json.loads(url_list[0])
                return content_have_url_dict

    def get_page_url(self):
        """
        用于获取每一页的url
        :return:
        """
        content_have_url_dict = self.homepage_script_dict
        html_urls = json.loads(content_have_url_dict['readerInfo']['htmlUrls'])
        if self.file_type == "ppt":
            return html_urls
        return html_urls['json']

    def get_file_type(self):
        """
        获取文件类型
        :return:
        """
        content_have_url_dict = self.homepage_script_dict
        return content_have_url_dict['viewBiz']['docInfo']['fileType']

    def get_page_content(self, page_url):
        """
        获取每一页的内容
        :return:
        """
        res = self.with_cookie(page_url)
        # 乱码变中文
        res_cn = res.text.encode('utf-8').decode('unicode_escape', 'ignore')
        res_cn_list = re.findall(r"\((.+)\)", res_cn)

        res_cn_dict = json.loads(format_dict_str(res_cn_list[0]))
        res_content_list = res_cn_dict['body']
        res_list = []

        for co_dict in res_content_list:
            res_list.append([co_dict['c'], co_dict['p']['y']])

        res_content_list = []

        for i, val in enumerate(res_list):
            if i != 0 and str(val[1]) == str(res_list[i - 1][1]):
                res_content_list[-1] = str(res_content_list[-1]) + str(val[0])
            else:
                res_content_list.append(val[0])

        return res_content_list

    def download(self, doc_type="txt"):
        if self.file_type == "ppt":
            self.download_in_ppt()
        elif self.file_type == "word":
            self.download_in_docx()
        else:
            self.download_in_txt()

    def download_in_ppt(self):
        url_info_dict_list = self.get_page_url()

        pictures = []
        save_dir = "./" + self.title.strip()
        try:
            os.mkdir(save_dir)
        except FileExistsError:
            print("文件夹已存在")
        pptx_name = '{}.{}'.format(save_dir + "/" + self.title, 'pptx')
        for jpg_num, jpg_url in enumerate(url_info_dict_list):
            response = requests.get(jpg_url)

            # 写入文本
            # print("开始下载第{}页数据".format(url_info_dict['pageIndex']))
            save_name = "/".join((save_dir, str(jpg_num))) + '.jpg'
            pictures.append(save_name)
            with open(save_name, "wb") as f:
                f.write(response.content)

        # 创建ppt
        ppt_file = pptx.Presentation()

        for fn in pictures:
            slide = ppt_file.slides.add_slide(ppt_file.slide_layouts[1])

            # 为PPTX文件当前幻灯片中第一个文本框设置文字，本文代码中可忽略
            # slide.shapes.placeholders[0].text = fn[:fn.rindex('.')]

            # 导入并为当前幻灯片添加图片，起始位置和尺寸可修改
            slide.shapes.add_picture(fn, Inches(0), Inches(0), Inches(10), Inches(7.5))

            ppt_file.save(pptx_name)

    def download_in_txt(self):
        url_info_dict_list = self.get_page_url()
        for url_info_dict in url_info_dict_list:
            page_content = self.get_page_content(url_info_dict['pageLoadUrl'])

            # 写入文本
            print("开始下载第{}页数据".format(url_info_dict['pageIndex']))
            if url_info_dict == url_info_dict_list[0]:
                open_model = 'w+'
            else:
                open_model = 'a+'
            with open("./" + self.title + ".txt", open_model, encoding='utf-8') as f:
                for content_li in page_content:
                    f.write(str(content_li) + "\n")
            f.close()

    def download_in_docx(self):
        # 创建文档对象
        document = Document()

        # 设置文档标题，中文要用unicode字符串
        # document.add_heading(self.title, 0)
        url_info_dict_list = self.get_page_url()
        for url_info_dict in url_info_dict_list:
            page_content = self.get_page_content(url_info_dict['pageLoadUrl'])
            print("开始下载第{}页数据".format(url_info_dict['pageIndex']))
            for content_li in page_content:
                # 添加无序列表
                paragraph = document.add_paragraph(str(content_li))
                paragraph_format = paragraph.paragraph_format
                paragraph_format.space_before = Pt(0)
                paragraph_format.space_after = Pt(0)
                paragraph.paragraph_format.line_spacing_rule = WD_LINE_SPACING.AT_LEAST
                paragraph_format.line_spacing = Pt(13)
        document.save("./" + self.title + ".docx")


if __name__ == '__main__':
    # 填写你的cookie，不填只能爬取部分
    cookie_str = ""
    # 百度文库url
    bidu_url = "https://wenku.baidu.com/view/bc18df72ba4ae45c3b3567ec102de2bd9605de7d.html"
    # main()
    bidu = BDWenKu(bidu_url, cookie_str)
    bidu.download()
